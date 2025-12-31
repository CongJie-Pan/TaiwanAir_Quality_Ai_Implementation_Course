# AIp08關連的挖掘X.py: AI python 實作 - 08: 關連的挖掘
# Jia-Sheng Heh, 10/22/2024, revised from AIp07聚類及評估.py
# Usage: streamlit run AIp08關連的挖掘X.py --> http:/localhost:8501

import numpy as np
import pandas as pd
# from st_aggrid import AgGrid #, GridUpdateMode, JsCode, ColumnsAutoSizeMode
from datetime import date, datetime
import plotly_express as px
import plotly.graph_objs as go
from pptx import Presentation   
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os
wkDir = "c:/Users/jsheh/Desktop/postWorking/DIKW/AIp/";   os.chdir(wkDir);   print(os.getcwd())

#%%##== 聚類使用之系統函式庫
import matplotlib.pyplot as plt
from sklearn.preprocessing import StandardScaler
from scipy.cluster.hierarchy import dendrogram, linkage, fcluster
import matplotlib.font_manager as fm
from sklearn.metrics import silhouette_score

#%%####### (W).網站系統基本架構 ##########
import streamlit as st
from streamlit_navigation_bar import st_navbar

#%%##===== (W1).應用相關資料: 分析參數 + 應用函式庫 + streamlit快取機制 =====#####


##== (1).定義相關分析參數: 各家企業的以下參數不盡相同 ==##
FFbreaks = [0, 1, 9, 99, 999, 19999]
MMbreaks = [-5000, 0, 999, 9999, 99999, 999999, 19999999]
BBbreaks = [0, 1, 7, 30, 99, 300, 1999]
RRbreaks = [0, 7, 30, 60, 99, 180, 360, 499, 700, 1999]               #--> 用於客戶漏斗 (下述)
Tnow = pd.to_datetime("2017/12/31", format="%Y/%m/%d");  print(Tnow)  # -- 數據分析點: 2023-07-01 00:00:00

#%%== (2).應用函式庫 (Cv): buildCv(),NES3(),addCvNES3() ==##
def buildCv(XX,FFbreaks,MMbreaks,BBbreaks):  ##== 建構客戶價值數據框: Cv = buildCv(X,FFbreaks,MMbreaks,BBbreaks)
    Cv1 = XX.groupby("customer").agg({"invoiceNo": "nunique", "amount": "sum", "quantity": "sum",
                                      "date": ["min", "max", "nunique"],
                                      })
    Cv1.columns = ["FF", "MM", "TT", "D0", "Df", "DD"]
    Cv1["FF0"] = pd.cut(Cv1["FF"], bins=FFbreaks).astype(str)
    Cv1["MM0"] = pd.cut(Cv1["MM"], bins=MMbreaks).astype(str)
    Cv1["BB"] = [(Cv1["Df"][k]-Cv1["D0"][k]).days/(Cv1["FF"][k]-1)
                 if (Cv1["FF"][k] != 1) else np.nan for k in np.arange(Cv1.shape[0])]
    Cv1["BB0"] = pd.cut(Cv1["BB"], bins=BBbreaks).astype(str)
    Cv1["yq0"] = pd.PeriodIndex(Cv1.D0, freq='Q')
    Cv1["yqf"] = pd.PeriodIndex(Cv1.Df, freq='Q')
    #              FF     MM  TT         D0  ...            MM0     BB   yq0   yqf
    # customer                               ...
    # A0000000036   3  10488   6 2021-01-31  ...  (9999, 49999]  188.0  21Q1  22Q1
    # A0000000038   5  20735   8 2020-08-04  ...  (9999, 49999]   93.5  20Q3  21Q3
    return Cv1
def NES3(Ck, K, M):                          ##== status = NES3(Ck, K, M): 定義NES3狀態
    if Ck["R0"] < 0: status = "U尚未消費"
    elif Ck["R0"] < 2*K:
        if Ck["MM"] > M: status = "N1新貴客"
        else:            status = "N2新客"
    else:
        if Ck["Rf"] < 2*K:
            if Ck["R0"]/Ck["FF"] < 0.75*K: status = "A1較活躍客"
            else:                          status = "A2活躍客"
        elif Ck["Rf"] < 3*K: status = "S1瞌睡客"
        elif Ck["Rf"] < 4*K: status = "S2半睡客"
        elif Ck["FF"] < 10:  status = "S3沈睡客"
        else:                status = "S4沈睡忠誠客"    
    return status
def addCvNES3(Cv, Tnow, KK, MM, RRbreaks):   ##== 將 NES3 加入 Cv: Cv = addCvNES3(Cv, Tnow, KK, MM, RRbreaks)
    Cv["R0"] = [(Tnow - pd.to_datetime(d)).days for d in Cv.D0]
    Cv["Rf"] = [(Tnow - pd.to_datetime(d)).days for d in Cv.Df]
    Cv["R00"] = pd.cut(Cv["R0"], bins=RRbreaks).astype(str)
    Cv["Rf0"] = pd.cut(Cv["Rf"], bins=RRbreaks).astype(str)
    Cv["status"] = Cv.apply(NES3, K=KK, M=MM, axis=1)
    return Cv

#%%== (2A).hclust聚類函式庫 ==##
scaler = StandardScaler()
font_path = 'C:/Users/jsheh/Desktop/newWorking/RDsys/RDSgpt/微軟正黑體-1.ttf'   #-- # 加載微軟正黑體字型
def plotDendrogram(Z, font_path):    ##== 依據聚類好的linkage(Z)繪製樹圖(dendrogram),以font_path為中文字體 ==##
    prop = fm.FontProperties(fname=font_path)
    plt.rcParams['font.family'] = prop.get_name()       #-- 使用該字體進行繪圖
    # 繪製 Dendrogram 圖
    plt.figure(figsize=(10, 7));   plt.title("Dendrogram 用於聚類數判斷", fontproperties=prop)
    dendrogram(Z)
    plt.xlabel('客戶', fontproperties=prop);   plt.ylabel('距離', fontproperties=prop);   
    fig = plt.gcf();   plt.close()   #  plt.show()
    return fig
def plotSilhouette(Z,CP):            ##== 測試聚類linkage(Z<-CP) 2-10個聚類數，計算其輪廓係數(silhouette_scores)繪圖並求最佳聚類數(best_k) ==##
    silhouette_scores = []
    K_range = range(2, 11)
    for k in K_range:
        cluster_labels = fcluster(Z, k, criterion='maxclust')   #-- 使用 fcluster 來獲取聚類結果
        silhouette_avg = silhouette_score(CP, cluster_labels)
        silhouette_scores.append(silhouette_avg)
    ##-- 畫出聚類數與輪廓係數的關係圖
    plt.figure(figsize=(8, 5));                     plt.plot(K_range, silhouette_scores, 'bx-')
    plt.xlabel('聚類數量 (k)');                     plt.ylabel('輪廓係數 (Silhouette Score)')
    plt.title('層次聚類的聚類數量與輪廓係數的關係');   fig = plt.gcf();   plt.close()   # plt.show()
    ##-- 顯示最佳聚類數
    best_k = K_range[silhouette_scores.index(max(silhouette_scores))]
    print(f"最佳聚類數量為: {best_k}")
    return fig, best_k, silhouette_scores
def clusterFeature(CP,Z,n_clusters): ##== 求取CP矩陣之linkage Z, 以n_clusters為聚類, 求取聚類CP,各類客戶數customer_clusters,及各類品類特徵cluster_characteristics ==##
    CP['cluster'] = fcluster(Z, n_clusters, criterion='maxclust')   #-- 使用 fcluster 來根據聚類數為 6 進行聚類
    ##-- 查看每個類別中的客戶數量
    customer_clusters = CP.groupby('cluster').size()
    print("各類別的客戶數量：");   print(list(customer_clusters))     #-- [4, 356, 7, 7, 1, 1]
    ##-- 分析各類別的品類特徵
    cluster_characteristics = CP.groupby('cluster').sum()
    print("\n各類別的品類特徵：");   print(cluster_characteristics)   #-- 查看各類別的品類特徵
    # category  kind1  kind11  kind12  kind16  kind17  kind2  kind27  kind3  kind56  kind6
    # cluster                                                                             
    # 1             0       0       0       0       0      0       0      0     252      0
    # 2          3249     525    1456     658     811   5578     175   1375      58    129
    # 3           749      20     129       3      15     77      41     28       0     49
    # 4            93     103      59     273      57    666      23    176       0      8
    # 5            35      48     104      10       6    190       1    105       0      7
    # 6             5       9      35       0     405     60       0     14       0      0
    return CP, customer_clusters, cluster_characteristics
def generate_cluster_descriptions(cluster_characteristics, customer_clusters, top_n=3): ##== 由各類客戶數customer_clusters,及各類品類特徵cluster_characteristics, 求取各類描述descriptions ==##
    descriptions = []
    for cluster_num in customer_clusters.index:         ##== 遍歷每個類別
        row = cluster_characteristics.loc[cluster_num]       #-- 獲取該類別的特徵數據        
        customer_count = customer_clusters[cluster_num]      #-- 獲取該類別的客戶數
        avg_counts = (row / customer_count).sort_values(ascending=False)  #-- 計算平均每位客戶的消費數量，並轉換為 DataFrame 以便排序
        top_categories = avg_counts.head(top_n)              #-- 只取前 top_n 大的品類
        description_parts = []                               #-- 初始化描述列表
        for category, avg_count in top_categories.items():   #== 遍歷前 top_n 的品類
            if avg_count >= 10:  description_parts.append(f"數十件*{category}")
            elif avg_count >= 1: description_parts.append(f"{avg_count:.1f}件*{category}")
            else:                description_parts.append(f"少量*{category}")
        descriptions.append(" + ".join(description_parts))   #-- 將合併的描述加入列表
    return descriptions
def calculate_within_between_variance(data, cluster_labels):   ##== 計算數據點/類別中心 的類別內/間變異數, 及各類別的類別內變異數 ==##
    overall_mean = np.mean(data, axis=0)          #-- 計算總體數據的均值（整體中心）
    unique_clusters = np.unique(cluster_labels)   #-- 獲取唯一的類別標籤，並計算總的類別數
    num_clusters = len(unique_clusters)
    within_cluster_variance = 0                   #-- 初始化類別內變異數和類別間變異數
    between_cluster_variance = 0
    within_cluster_variances = [0] * num_clusters #-- 用來存儲每個類別的類別內變異數
    for cluster in unique_clusters:    #== 計算每個類別的均值
        cluster_data = data[cluster_labels == cluster]  #-- 提取當前類別的數據點
        cluster_mean = np.mean(cluster_data, axis=0)    #-- 計算該類別的均值（中心）
        cluster_within_variance = np.sum((cluster_data - cluster_mean) ** 2)  #-- 計算類別內變異數（每個點與其類別均值的距離平方和）
        within_cluster_variance += cluster_within_variance        
        # 使用 list 儲存每個類別的類別內變異數
        within_cluster_variances[cluster - 1] = cluster_within_variance  # 假設 cluster 是從 1 開始編號        
        between_cluster_variance += len(cluster_data) * np.sum((cluster_mean - overall_mean) ** 2)  #-- 計算類別間變異數（類別均值與總體均值的距離平方和）
    return within_cluster_variance, between_cluster_variance, within_cluster_variances
def assocRule(XX,min_supp,min_conf): ##== 求取關聯規則: assocRule(XX,min_supp,min_conf)
    CPlist = XX.groupby('customer')['category'].apply(lambda x: list(set(x.unique()))).tolist()  ##== (1)
    transactions = CPlist
    total_transactions = len(transactions)
    from mlxtend.preprocessing import TransactionEncoder                         ##== (2)
    te = TransactionEncoder()
    te_ary = te.fit(transactions).transform(transactions)
    oht_df = pd.DataFrame(te_ary, columns=te.columns_);   
    from mlxtend.frequent_patterns import apriori, association_rules
    FIS = apriori(oht_df, min_support=min_supp, use_colnames=True, max_len=5)    ##== (3): min_supp = 0.4
    FIS['count'] = FIS['support'] * len(oht_df)
    FIS.sort_values(by='count', ascending=False, inplace=True)
    FIS['count'] = FIS['support'] * len(oht_df)
    rules = association_rules(FIS, metric="confidence", min_threshold=min_conf)  ##== (4): min_conf = 0.6
    rules['count'] = (rules['antecedent support'] * total_transactions).astype(int)
    rules = rules[(rules['antecedents'].apply(len) == 1) & (rules['consequents'].apply(len) == 1)]  # <== 去除掉 2-項集等
    # rules['antecedents'] = rules['antecedents'].apply(lambda x: ', '.join(list(x)))
    # rules['consequents'] = rules['consequents'].apply(lambda x: ', '.join(list(x)))
    return rules
def graph_style(rules):              ##== 從關聯規則(rules,及選取數據): graph, style = graph_style(rules,XX)
    import igraph as ig
    graph_edges = [(list(rule['antecedents'])[0], list(rule['consequents'])[0]) for _, rule in rules.iterrows()]
    g = ig.Graph.TupleList(graph_edges, directed=True);   print(g)    ##== (1)
    lenXX = rules["count"][0]/rules["antecedent support"][0]
    node_size_dict = {vertex: rules[rules['antecedents']=={vertex}]['antecedent support'].iloc[0]*lenXX 
                      if not rules[rules['antecedents']=={vertex}].empty else 1 for vertex in g.vs['name']}  ##== (2)
    node_labels = [f"{vertex['name']}\n ({int(node_size_dict[vertex['name']] )})" for vertex in g.vs]
    node_sizes  = [ 0.5*node_size_dict.get(vertex['name'],1) for vertex in g.vs];  
    print(node_sizes)  #-- [169.0, 129.0, 124.5, 118.5, 90.0] ---> 實驗調整節點大小
    visual_style = { "layout": g.layout("fr"), "bbox": (900, 700), "margin": 50,
        "vertex_size": node_sizes, "vertex_label": node_labels, 'vertex_color': 'orange',
        "edge_label": [f"{rule['confidence']:.2f}" for _, rule in rules.iterrows()], 'edge_color': 'lightblue',
        "edge_arrow_size": 2, }
    return g, visual_style

#%%== (3).數據st.cache函式庫: getX(),buildCvRDS() ==##
@st.cache_data
def getX(Xname):     ##== X = getX(Xname): 自 X.csv 讀取 X (KDD1), 並設定標籤 (KDD3) ==##
    X = pd.read_csv(Xname)
    # -- 還有很多其他產生此標籤的方法, 這裡只是取其中較方便的一種
    X["date"] = pd.to_datetime(X["datetime"]).dt.date
    X["year"] = pd.to_datetime(X["datetime"]).dt.year
    X["yq"] = pd.PeriodIndex(X.date, freq='Q')
    X["ym"] = pd.PeriodIndex(X.date, freq='M')
    return(X)
@st.cache_data
def buildCvRDS(X,FFbreaks,MMbreaks,BBbreaks,RRbreaks,Tnow):  ##== Cv = buildCvRDS(X,..): 由交易數據 X 求取客戶數據框 Cv (KDD3)
    Cv = buildCv(X,FFbreaks,MMbreaks,BBbreaks)
    print(Cv.shape);   print(Cv[2:4])  # -- (52217, 17)
    KK = np.nanmean(Cv["BB"]);    print(KK)  # -- 43.070694784611675
    MM = np.nanmean(Cv["MM"]);    print(MM)  # -- 46998.990443725226
    Cv = addCvNES3(Cv, Tnow, KK, MM, RRbreaks);   print(Cv.shape);   print(Cv[2:4])  # -- (52217, 22)
    # Cv = pd.read_csv("cvv.csv")
    return(Cv)

#%%##===== (W2).儀表板函式庫: 前台(a)navbar, 後台(b), (c)canvas,(d)sidebar, (e)返回值,(f)LOG =====#####
def 擷取交易(Xname):      ##== (KDD1) X, rMsg, log = 擷取交易(Xname) ==##
    ##== (b).後台 ==##
    X = getX(Xname);    print("\n\n>>>>> 擷取交易數據 (-->XXX) -----")  # -- 偵錯用
    ##== (c).前台-canvas ==##
    st.header("== (KDD1).X=擷取交易(Xname) -- ")     
    st.write(f"* 交易數據檔 Xname ={Xname}" )
    st.write(f"* 交易數據框 X 具有 {X.shape[0]} 筆記錄")
    st.subheader("(KDD1).交易數據框 X -- ")
    st.dataframe(X.head(3))
    # == (d).前台-sidebar ==##  #==> [[AIp04/C4)(2)垂直流程]]
    st.sidebar.header("== (KDD1).X=擷取交易(Xname) -- ")
    st.sidebar.write(f"* 交易數據檔 Xname ={Xname}" )
    st.sidebar.write(f"* 交易數據框 X 具有 {X.shape[0]} 筆記錄")
    ##== (e/f).返回信息(rMsg)/LOG(log) ==##
    rMsg = [ { "Ptitle": f"擷取交易({Xname})", 
               "Plist": ["(KDD1).擷取交易數據 -- ",f"交易檔名 = {X.shape[0]}", f"記錄筆數 = {X.shape[0]}筆" ],
               "Plevel": [ 0, 1, 1 ], "df": ["X"] } ]
    log = ["[[操作]] X=擷取交易(Xname)", f"[參數] 交易數據檔 Xname ={Xname} ", f"[結果] 交易數據框 X 具有 {X.shape[0]} 筆記錄"]
    return X, rMsg, log
def 轉換客戶圖像(X,FFbreaks,MMbreaks,BBbreaks,RRbreaks,Tnow): ##== (KDD3,4) Cv, TFM, rMsg, log = 轉換客戶圖像(X,FFbreaks,MMbreaks,BBbreaks,RRbreaks,Tnow) ==##
    ##== (a).後台 ==##
    Cv = buildCvRDS(X,FFbreaks,MMbreaks,BBbreaks,RRbreaks,Tnow)
    # TFM = pd.crosstab(Cv["FF0"], Cv["MM0"], margins=True);   print(TFM)
    # TFMs = pd.DataFrame(TFM);     print(TFMs)  #-- (B).客戶價值模型
    # TFMs_styled = TFMs.style.format(formatter="{:,}", na_rep=".").bar(cmap="cool", axis=None)  #==> 表格加上style.format更豐富
    FFdata = pd.DataFrame(Cv.FF0.value_counts()).reset_index()              #-- (C1).客戶FF漏斗   
    figFF = px.funnel(x=list(FFdata.FF0), y=list(FFdata.index.astype(str)), title="造訪頻次客戶漏斗") #==> px有漏斗的元件
    MMdata = pd.DataFrame(Cv.MM0.value_counts()).reset_index().sort_values(by='MM0');   print(MMdata) #-- (C2).客戶MM漏斗
    figMM = px.funnel(x=list(MMdata.MM0), y=list(MMdata["MM0"].astype(str)), title="消費金額客戶漏斗")
    ##== (c1).前台-canvas: 客戶圖像 ==##
    st.header("== (KDD3).客戶圖像(Cv) --")
    with st.expander("___ 由交易數據(X)可以轉換出客戶圖像(Cv), 其樣貌如下 ....."): 
        st.subheader("* (A) 共有" + str(Cv.shape[0]) + "位客戶")
        st.subheader("* (B) (KDD4) 客戶數據框(Cv)")
        st.dataframe(Cv.head(2))
    # st.table(data=TFMs_styled)   #-- st.dataframe(TFM)
    ##== (c2).前台-canvas: 客戶漏斗 (Cv.FF0-->FFdata, Cv.MM0-->MMdata) ==##
    st.header("== (KDD4).客戶漏斗 --")
    with st.expander("___ 客戶圖像中的造訪次數區間(FF0)與消費金額區間(MM0), 可以下列圖示表現 ....."): 
        cols = st.columns([1, 1])
        cols[0].subheader("* (C1) 造訪次數(FF)客戶漏斗");   cols[0].plotly_chart(figFF, theme="streamlit", use_container_width=True)
        cols[1].subheader("* (C2) 消費金額(MM)客戶漏斗");   cols[1].plotly_chart(figMM, theme="streamlit", use_container_width=True)
    ##== (d).前台-sidebar ==##  
    st.sidebar.header("== (KDD3).客戶圖像 --")
    st.sidebar.write("* 共有" + str(Cv.shape[0]) + "位客戶")
    ##== (e/f).返回信息(rMsg)/LOG(log) ==##
    figFF.write_image("figFF.jpg")
    figMM.write_image("figMM.jpg")
    rMsg = [ { "Ptitle": "Cv = 轉換客戶圖像(X,..)", 
               "Plist":  ["(KDD2).參數設定 -- ", f"FFbreaks = {FFbreaks}", f"MMbreaks = {MMbreaks}", f"BBbreaks = {BBbreaks}" ],
               "Plevel": [ 0, 1, 1, 1 ] },
             { "Ptitle": "(KDD3).客戶圖像(Cv)", 
               "Plist":  ["(KDD3).客戶圖像 --", f"共有{Cv.shape[0]}位客戶", "(KDD3)客戶數據框(Cv)" ],
               "Plevel": [ 0, 1, 1 ], "df": ["Cv"] },
             { "Ptitle": "(KDD4).客戶漏斗", 
               "Plist":  ["(KDD4).客戶漏斗 --", "造訪次數(FF)客戶漏斗", "消費金額(MM)客戶漏斗" ],
               "Plevel": [ 0, 1, 1 ], "fig": ["figFF.jpg","figMM.jpg"] },
             ]
    log = ["[[操作]] Cv = 轉換客戶圖像(Xname,FFbreaks,MMbreaks,BBbreaks,RRbreaks,Tnow)",
           f"[參數] 造訪頻次區間 FFbreaks = {FFbreaks}", f"[參數] 消費金額區間 MMbreaks = {MMbreaks}", 
           f"[參數] 回購週期區間 BBbreaks = {BBbreaks}", f"[參數] 最後造訪區間 RRbreaks = {RRbreaks}" ,
           f"[[結果]] 客戶圖像 Cv 具有 {Cv.shape[0]} 位客戶",
            "[[結果]] Cv.FF 與 Cv.MM 可用以繪製客戶漏斗", ]
    return Cv, rMsg, log
def TFM_客群選取(Cv):     ##== (KDD4,5) TFMs, FF0A, MM0A, CvTA, rMsg, log = TFM_客群選取(Cv) ==##
    ##== (b).後台: 設定 TFMs, TFMvs, TFMv0s 之style ==##
    TFM = pd.crosstab(Cv["FF0"], Cv["MM0"], margins=True);   print(TFM)
    TFMs = pd.DataFrame(TFM);     print(TFMs)  #-- (B).客戶價值模型
    TFMs_styled = TFMs.style.format(formatter="{:,}", na_rep=".").bar(cmap="cool", axis=None)  #==> 表格加上style.format更豐富
    TFMv = pd.crosstab(index=Cv.FF0, columns=Cv.MM0,values=Cv.MM, aggfunc='sum', margins=True)
    TFMvs = pd.DataFrame(TFMv);   print(TFMvs)
    TFMvs_styled = TFMvs.style.format(formatter="{:,}", na_rep=".").bar(cmap="Wistia", axis=None)
    TFMv0s = pd.DataFrame(100*TFMv/TFMv["All"]["All"]);    print(TFMv0s)
    TFMv0s_styled = TFMv0s.style.format(formatter="{:,.1f}", na_rep=".").bar(cmap="Wistia", axis=None)
    #== (c).前台-canvas: 客戶價值模型 ==##
    st.header("== (KDD4).客戶價值模型 --")
    with st.expander("___ Cv中的造訪頻次(FF0)與消費金額(MM0), 可以得到價值模型 ....."): 
        st.subheader("* (A) 共有" + str(Cv.shape[0]) + "位客戶")
        st.subheader("* (B) 客戶價值模型")
        st.table(data=TFMs_styled)   # st.dataframe(TFM)
    cols = st.columns([3, 2])    #== 營業額模型 (TFMvs,TFMv0s)
    cols[0].subheader("- (C1) 客戶價值營業額模型");       
    with cols[0]: 
        with st.expander("___ 展開....."): st.table(data=TFMvs_styled)    # st.dataframe(TFM)
    cols[1].subheader("- (C2) 客戶價值營業額佔比模型");   
    with cols[1]:
        with st.expander("___ 展開....."): st.table(data=TFMv0s_styled)   # st.dataframe(TFM)
    ##== (d).前台-sidebar ==##  
    ##-- (d1).客群選擇
    CvTA = Cv
    st.sidebar.header("== (KDD4).客戶價值模型 --")
    st.sidebar.subheader("- (D) 選擇客群,以觀看客戶圖像/旅程")
    FF0A = st.sidebar.multiselect('>> 請選擇客戶的造訪頻次區間(FF0):', list(Cv.FF0.unique()))
    FF0 = [str(x) for x in FF0A];     print(">>> "+str(FF0))
    MM0A = st.sidebar.multiselect('>> 請選擇客戶的消費金額區間(MM0):', list(Cv.MM0.unique()))
    MM0 = [str(x) for x in MM0A];     print(">>> "+str(MM0))
    ##-- (d2).客群計算
    cIND1 = np.arange(Cv.shape[0]);   cIND = []
    if len(FF0) > 0:  cIND1 = np.where([(str(Cv.FF0[k]) in FF0) for k in cIND1])[0]
    if len(MM0) > 0:  cIND = np.where([(str(Cv.MM0[k]) in MM0) for k in cIND1])[0]
    st.sidebar.write("-> 目標客戶數 = "+str(len(cIND))+"位")
    if len(cIND)!=0:
        CvTA = Cv.iloc[cIND];   
        # indX = list(np.where([k if X.customer[k] in CvTA.index else None for k in np.arange(X.shape[0])])[0])  # indX
        CvTA.reset_index(inplace=True)
        st.sidebar.write("-> 涵蓋交易數 = "+str(sum(CvTA.FF))+"筆")
    ##-- (d3).客戶圖像+客戶旅程
    flagCvTA = False
    if st.sidebar.checkbox("- (E).確定客群"):
        setCvTA = "客戶價值 FF0="+",".join(FF0A) +", MM0A="+",".join(MM0A)+", 共計"+str(len(CvTA))+"位";   
        st.subheader("* (E1) 目標客戶: FF0 = "+",".join(FF0) +", MM0 = "+",".join(MM0)+", 共計"+str(len(CvTA))+"位")
        st.subheader("* (E2) 客戶圖像")
        st.dataframe(CvTA)
        ##-- (d4).匯出至EXCEL檔
        if st.sidebar.checkbox("* (E) 匯出客群名單 至EXCEL檔"):
            CvTA.to_excel(setCvTA+".xlsx" );   flagCvTA = True       #==> pandas 提供匯出 EXCEL 的方法 .to_excel()
            st.sidebar.write("* (E)"+str(len(cIND))+f"筆記錄已輸出至檔案{setCvTA}.xlsx")
    ##== (e/f).返回信息(rMsg)/LOG(log) ==##   
    rMsg = [ { "Ptitle": "(KDD4).客戶價值模型", 
               "Plist":  ["客戶價值模型 -- ", f"共有 {Cv.shape[0]} 位客戶", f"總消費金額為 {np.sum(Cv.MM)} 元", "客戶價值模型" ],
               "Plevel": [ 0, 1, 1, 0 ], "table": ["TFMs"] },
             { "Ptitle": "TFM模型所選取到的目標客群(CvTAtfm)", 
                         "Plist":  ["(KDD4).目標客群 -- ", f"TFM參數: FF0A={FF0A}, MM0A={MM0A}", f"目標客群人數 = {CvTA.shape[0]} 位" ],
                         "Plevel": [ 0, 1, 1 ], "df": ["CvTAtfm"] }, ]
    log = ["[[操作]] TFMs, FF0A, MM0A, CvTA = TFM_客群選取(Cv)", f"[結果] 客戶價值模型 TFMs 具有 {TFMs.shape} 個元素",
           f"[[選擇]] 造訪頻次區間 FF0A = {FF0A}", f"[[選擇]] 消費金額區間 MM0A = {MM0A}",    f"[結果] 目標客戶數據框 CvTAtfm 具有 {CvTA.shape[0]} 位客戶",]
    if flagCvTA: 
        rMsg.append( {"Ptitle": "匯出目標客群檔", "Plist": [ "匯出XLSX檔案", f"記錄數 = {CvTA.shape[0]} 位", f"檔名 = {setCvTA}.xlsx" ],"Plevel": [0,1,1] } )
        log.append(f"[結果] 匯出 CvTAtfm 為 {setCvTA}.xlsx 檔")
    return TFMs, FF0A, MM0A, CvTA, rMsg, log
def 目標客群交易(CvTA,X): ##== (KDD3,4) XX, rMsg, log = 目標客群交易(CvTA,X) ==##
    ##== (a).後台 ==##
    # print(">>>>>>> 目標客群交易: ");  print(sss)
    XX = X.loc[X["customer"].isin(CvTA["customer"])];     print(XX.shape);     print(XX.head(2))     #-- (19228, 17)   
    ##== (c).前台-canvas ==##
    st.header("== (KDD3).XX=目標客戶交易(CvTA,X) -- ")     
    st.write(f"* 目標客戶交易 XX 具有 {XX.invoiceNo.nunique()} 筆交易")
    st.write(f"* 目標客戶交易 XX 具有 {XX.shape[0]} 筆記錄")
    st.subheader("(KDD3).目標客戶交易數據框 XX -- ")
    st.dataframe(XX.head(3))
    # == (d).前台-sidebar ==##  #==> [[AIp04/C4)(2)垂直流程]]
    st.sidebar.header("== (KDD3).XX=目標客戶交易(CvTA,X) -- ")
    st.sidebar.write(f"* 目標客戶交易 XX 具有 {XX.shape[0]} 筆記錄")
    ##== (e/f).返回信息(rMsg)/LOG(log) ==##
    rMsg = [ { "Ptitle": "XX=目標客戶交易(CvTA,X)", 
               "Plist": ["(KDD3).提取目標客戶交易數據 -- ", f"記錄筆數 = {XX.shape[0]}筆" ],
               "Plevel": [ 0, 1 ], "df": ["XX"] } ]
    log = ["[[操作]] XX=目標客戶交易(CvTA,X)", f"[結果] 目標客戶交易 XX 具有 {XX.shape[0]} 筆記錄"]
    return XX, rMsg, log
def 客群聚類分析(XX):     ##== (KDD3,4) XX, rMsg, log = 客群聚類分析(XX) ==##
    import matplotlib.pyplot as plt
    from sklearn.preprocessing import StandardScaler
    from scipy.cluster.hierarchy import dendrogram, linkage, fcluster
    import matplotlib.font_manager as fm
    from sklearn.metrics import silhouette_score
    ##== (a).後台 ==##
    # print(">>>>>>> 目標客群交易: ");  print(sss)
    XX_clean       = XX   #-- XX.drop(columns=["Unnamed: 0"])   #-- 移除無用的
    top_categories = XX_clean['category'].value_counts().nlargest(10).index   #-- 選取前十大品類
    filtered_XX    = XX_clean[XX_clean['category'].isin(top_categories)];   print(filtered_XX.shape);   print(filtered_XX.head(2))   #-- (17869, 17)
    CP_matrix = pd.crosstab(filtered_XX['customer'], filtered_XX['category']);   print(CP_matrix.shape);   print(CP_matrix.head(2))  #-- (376, 10) 客戶-品類的交叉表
    ##== (c).前台-canvas ==##
    st.header("== (KDD4).目標客戶聚類(XX) -- ")     
    st.write(f"* (1).目標客戶交易 XX 具有 {XX.shape[0]} 筆記錄")
    st.subheader("(KDD3).客戶品類交叉表 (CP_matrix) -- ")
    st.write(f"* (2).客戶品類交叉表 {CP_matrix.shape[0]} 個元素")
    st.dataframe(CP_matrix.head(3))
    ##== (d).前台-sidebar ==##  #==> [[AIp04/C4)(2)垂直流程]]
    st.sidebar.header("== (KDD4).目標客戶聚類(XX) -- ")
    st.sidebar.write(f"* 目標客戶交易 XX 具有 {XX.shape[0]} 筆記錄")
    st.sidebar.write(f"* 客戶品類交叉表 {CP_matrix.shape} 個元素")
    ##==
    if st.sidebar.checkbox("繪製聚類結果樹"):        
        CP_matrix_scaled = scaler.fit_transform(CP_matrix);   print(np.round(CP_matrix_scaled[0:2],2))  #-- 標準化數據
        # [[ 0.29  0.03 -0.45 -0.28 -0.15 -0.29 -0.37 -0.41 -0.12  1.85 -0.16]
        #  [-0.49 -0.42 -0.57  0.61  0.02 -0.25  0.21 -0.53 -0.12 -0.38 -0.16]]
        Z = linkage(CP_matrix_scaled, method='ward')        #-- 使用層次聚類中的 linkage 方法，使用 'ward' 方法進行聚類
        plt1 = plotDendrogram(Z,font_path)
        plt2, best_k, silhouette_scores = plotSilhouette(Z,CP_matrix_scaled);  # print(best_k)   #-- 2
        st.header("== (KDD4).客戶聚類分析(CP_matrix) -- ")     
        cols = st.columns([1, 1])    #== 
        cols[0].subheader("- (1) 聚類樹圖:");          cols[0].pyplot(plt1)
        cols[1].subheader("- (2) 輪廓係數圖:");        cols[1].pyplot(plt2);    cols[1].write(f"* 最佳聚類數 = {best_k}");   
        st.subheader( f"- (3) Silhouette Score: {np.round(silhouette_scores,2)}" )
        Ncls = st.sidebar.slider('- 聚類數 = ', 0, 20, best_k)
        if st.sidebar.checkbox(f"確認聚類數目={Ncls}"):     
            CP_matrixC, clusterSize, clusterFeatures = clusterFeature(CP_matrix,Z,n_clusters=Ncls);   #<-- 硬取聚類數=6
            cluster_descriptions = generate_cluster_descriptions(clusterFeatures, clusterSize);   print(cluster_descriptions)
            wcv, bcv, pwcv = calculate_within_between_variance(scaler.fit_transform(CP_matrix[CP_matrix.columns[0:10]]), CP_matrix["cluster"].to_numpy());    
            CQ = bcv / wcv
            AAA = clusterFeatures;   AAA["size"] = clusterSize;  AAA["類別內變異數"] = pwcv 
            AAA["description"] = cluster_descriptions;    print(AAA)
            st.header("== (KDD4-5).客戶聚類結果與解讀(cluster...) -- ")     
            st.subheader("- cluster + Features + Sizes + Descriptions:")
            st.dataframe(AAA)
            kk = st.sidebar.slider('請選擇一個數字', 1, Ncls, 1)
            st.header("== (KDD5).聚類評估 -- ")     
            st.subheader(f"(1) 類別內變異數 = {wcv:.3f}");   
            st.subheader(f"(2) 類別間變異數 = {bcv:.3f}");   
            st.subheader(f"(3) 聚類質量指標 = {CQ:.3f}");   
            st.header(f"== (KDD5).觀看第{kk}類 -- ")     
            st.dataframe(CP_matrixC.loc[CP_matrixC["cluster"]==kk][CP_matrixC.columns[0:10]])
    ##== (e/f).返回信息(rMsg)/LOG(log) ==##
    rMsg = [ { "Ptitle": "(KDD4).目標客戶聚類(XX)", 
               "Plist": ["(KDD3).目標客戶交易數據 -- ", f"* 目標客戶交易 XX 具有 {XX.shape[0]} 筆記錄", 
                         "(KDD3).客戶品類交叉表 (CP_matrix) -- ", f"* 客戶品類交叉表 {CP_matrix.shape[0]} 個元素" ],
               "Plevel": [ 0, 1, 0, 1 ], "df": ["CP_matrix"] } ]
    log = ["[[操作]] CP_matrix=目標客戶聚類(XX)", f"* 目標客戶交易 XX 具有 {XX.shape[0]} 筆記錄", f"* 客戶品類交叉表 {CP_matrix.shape} 個元素"]
    return XX, rMsg, log
def 客群關連規則(XX):     ##== (KDD3,4) XX, rMsg, log = 客群關連規則(XX) ==##
    from PIL import Image
    import igraph as ig
    ##== (a).後台 ==##
    CPlist = XX.groupby('customer')['category'].apply(lambda x: list(set(x.unique()))).tolist()  ##== (1)
    ##== (c).前台-canvas ==##
    st.header("== (KDD3-4).客戶關連分析(XX) -- ")     
    st.subheader("(1).(KDD3) 目標客群之交易 (XX) -- ")
    st.write(f"* 目標客群交易 XX 具有 {XX.shape[0]} 筆記錄")
    st.subheader("(2).(KDD3) 每位客戶之交易 (CPlist) -- ")
    st.write(f"* 目標客群有 {len(CPlist)} 位客戶, 前三筆為...")
    for i, sublist in enumerate(CPlist[:5]): st.markdown(f"- 第{i}位客戶的購買品類: {', '.join(sublist)}")
    # st.json( CPlist[0:5] )
    ##== (d).前台-sidebar ==##  #==> [[AIp04/C4)(2)垂直流程]]
    st.sidebar.header( "== (1-2).(KDD3) 客戶關連分析之數據準備(XX) -- " )
    st.sidebar.write(f"* 目標客群有 {len(CPlist)} 位客戶")
    st.sidebar.subheader("== (3).(KDD4) 客戶關連分析(XX) -- ")     
    minSupp = st.sidebar.slider('- 支持度 = ', min_value=0.001, max_value=0.99, value=0.4, step=0.05)
    minConf = st.sidebar.slider('- 置信度 = ', min_value=0.1,   max_value=0.99, value=0.6, step=0.1)
    if st.sidebar.checkbox(f"確認關連參數 minSupp={minSupp}, minConf={minConf}"):                 
        rules = assocRule(XX, minSupp, minConf);   print(rules[0:2])   #--> 答案如上方 (rules) --> 多了count欄位 
        graph, style = graph_style(rules)
        figASSOC = ig.plot(graph, **style).save('igraph_plot.png')
        graph_image = Image.open('igraph_plot.png')
        st.sidebar.write(f"* 求得的關連規則有 {rules.shape[0]} 條規則:")                
        st.subheader("- (4) 關連規則:");          
        st.write(f"* 求得的關連規則有 {rules.shape[0]} 條規則:")                
        st.write(f"* 求取關連規則之參數為 minSupp={minSupp}, minConf={minConf}")                
        st.dataframe(rules)
        st.subheader("- (5) 關連規則圖:"); 
        st.image(graph_image, use_column_width=True)
        # cols[1].plotly_chart(figASSOC, theme="streamlit", use_container_width=True)
    ##== (e/f).返回信息(rMsg)/LOG(log) ==##
    rMsg = [ ]
    log = [ ]
    return XX, rMsg, log
def 匯出PPT檔(PPTname):   ##== (KDD5) 匯出PPT檔(PPTname) ==##
    print(sss)
    #== (c1/d1).標題: 前台-canvas/前台-sidebar ==##
    st.header("(KDD5).匯出PPT檔案")
    st.sidebar.header("(KDD5).匯出PPT檔案");     
    if st.sidebar.checkbox("* (KDD5) 匯出至 PPTX檔"):
        #== (b).後台-存檔 ==##
        sug_list = [ { "Ptitle": "謝謝", "Plist":  ["~~敬請指教!!"], "Plevel": [ 0 ] } ];    appendPPTX(sug_list)
        sss.prs.save(PPTname)
        #== (c2/d2).前台-canvas: 匯出 ==##
        st.subheader(f"* (KDD5) 匯出LOG 至PPT檔--{PPTname} --")
        st.write(f"* (KDD5) PPT檔--{PPTname} 己匯出")
        st.sidebar.write("* (KDD5) PPTX 己匯出")
        #== (c2/d2).前台-canvas: 匯出 ==##
        st.subheader("* PPT檔案結構 --")
        for slide_number, slide in enumerate(sss.prs.slides, start=1):
            title = None
            for shape in slide.shapes:
                if shape.has_text_frame:  title = shape.text
            if title: st.write(f"* [第{slide_number}頁] {title}")
            else:     st.write(f"* [第{slide_number}頁] 沒有標題")
    return

#%%##===== (W3).本系統函式庫: (1)列印PPT,(2)導航 =====#####

#%%== (1).列印PPTX的函式庫: genPPTX(),addBulletPage(),addSlideDF(),makeDFtable() ==##
def genPPTX(mainTitle,subTitle):             ##== prs = genPPTX(mainTitle,subTitle): 產生一份新的投影片  
    from pptx import Presentation
    prs = Presentation()                 
    slide0 = prs.slides.add_slide(prs.slide_layouts[0])               
    slide0.shapes.title.text = mainTitle;    slide0.placeholders[1].text = subTitle
    return prs
def addBulletPage(prs,Ptitle,Plist,Plevel):  ##== prs = addBulletPage(prs,Ptitle,Plist,Plevel): 增加一個重點(Plist)頁,並設定重點層級(Plevel)及顏色 (Plevel=1)
    slide = prs.slides.add_slide( prs.slide_layouts[1] );  #-- 產生一頁(slide)新的 "標題與內容" 的重點頁(BulletPage)  
    slide.shapes.title.text = Ptitle                       #-- 設定標題(Ptitle)
    tf = slide.shapes.placeholders[1].text_frame           #-- 設定內文 文字框(tf)
    for k in np.arange(len(Plist)):
        if k==0:
            tf.text = Plist[0]   #-- 設定第 1 子標題 (tf.text = Plist[0])
        else:                    #-- 設定新增 子標題 (Plist[k]), 其層級 (Plevel[k]) 及顏色 (Plevel=1為粗體彩色)
            p = tf.add_paragraph();    
            p.level = Plevel[k];     p.text = Plist[k]   
            if (p.level==1): 
                p.font.bold = True
                p.font.color.rgb = RGBColor(0,0,255)  # RGBColor(0xFF, 0x7F, 0x50)
    print("addBulletPage>>> generate Bullet Page-"+Ptitle)        
    return prs
def addSlideDF(prs,ind,Ptable):              ##== prs = addSlideDF(prs,ind,Ptable): 將表格(Ptable)加入某頁 (prs.slides[ind])
    shapes = prs.slides[ind].shapes
    if (Ptable is not None):
        print("addSlideDF>>> generate dataframe Table...")        
        left, top, width, height = Inches(1), Inches(1), Inches(8), Inches(6)
        table = shapes.add_table(Ptable.shape[0], Ptable.shape[1], left, top, width, height).table        
        for i in np.arange(Ptable.shape[0]):
            for j in np.arange(Ptable.shape[1]):
                table.cell(i,j).text = str(list(Ptable.iloc[i])[j])
    return prs
def makeDFtable(df):                         ##== table = makeDFtable(df): make df to table with first row as column names
    Xcol = pd.DataFrame(df.columns).transpose();   Xcol.columns = df.columns;   
    AAA  = pd.concat([Xcol,df],axis=0); 
    Arow = pd.DataFrame(AAA.index);                Arow.index = Arow[Arow.columns[0]];   
    BBB  = pd.concat([Arow,AAA],axis=1);           BBB.index = Arow[Arow.columns[0]]
    return BBB

#%%== (2).導航函式庫: appendPPTX(),check2log(),initSSS() ==##
def appendPPTX(rv_list):              ##== 以傳回數據(rv_list)生成投影片sss.prs ==##
    print(rv_list)
    with st.sidebar.expander("___ 當前投影片生成步驟 ....."):    
        for i, rv in enumerate(rv_list):
            st.write(">> 生成投影片第"+str(len(sss.prs.slides)-1)+"頁-"+rv.get("Ptitle")+" 中...")
            # st.sidebar.write(">> 生成"+rv.get("Ptitle")+"投影片中...")
            print(rv)
            if rv.get("Ptitle"): sss.prs = addBulletPage(sss.prs, rv.get("Ptitle"), rv.get("Plist"), rv.get("Plevel"))
            # if rv.get("df"):  sss.prs = addSlideDF(sss.prs, len(sss.prs.slides)-1, sss[rv.get("df")].head(2))
            if rv.get("df"):
                
                dfList = rv.get("df");  
                for dfFile in dfList:  sss.prs = addSlideDF(sss.prs, len(sss.prs.slides)-1, makeDFtable(sss[dfFile].head(2)))
            if rv.get("table"):  
                tblList = rv.get("table");  
                for tbl in tblList:  sss.prs = addSlideDF(sss.prs, len(sss.prs.slides)-1, makeDFtable(sss[tbl]))
            if rv.get("fig"):    
                slide1 = sss.prs.slides[len(sss.prs.slides)-1];      picList = rv.get("fig");   
                for picFile in picList:  pic1 = slide1.shapes.add_picture(picFile, Inches(1), Inches(1));  #  print(">>>>> 2."+picFile)               
    return               
def check2log(strList,log,rv_list):   ##== check 再將 textStr 納入 log 中, 並中並可以提供建議, 並生成投影片 ==##
    st.sidebar.markdown('---') 
    st.sidebar.header("== (KDD5)數據解讀/LOG操作 --")
    with st.sidebar.expander("___ 用戶可以在下方直接輸入[用戶名]與[建議]進行鹬讀 .....", expanded=True): 
        ##-- 生成建議 --## 
        st.session_state.username = st.text_input("", st.session_state.username, placeholder="輸入用戶名")
        sugg_key = f"log_checkbox_{len(log)}"
        suggestion = st.text_area("", value=st.session_state.suggestion, key=sugg_key, height=20, placeholder="輸入建議")
        if st.button("LOG操作 / 提交建議"):
            ##-- 生成投影片 --##
            appendPPTX(rv_list)
            ##-- LOG/提交 --##
            if st.session_state.username and suggestion:
                log.extend(strList)
                log.append(f"<<{st.session_state.username}建議>> {suggestion}")
                sug_list = [ { "Ptitle": f"<<{st.session_state.username}建議>> {suggestion}", 
                               "Plist":  ["(KDD5).用戶建議 -- ", f"用戶 = {st.session_state.username}", f"建議 = {suggestion}" ],
                               "Plevel": [ 0, 1, 1 ] } ]
                appendPPTX(sug_list)
                st.success("建議已提交纳入LOG");     st.session_state.suggestion = "" 
            else:
                log.extend(strList)
                st.error("不列入建議,只單純LOG操作內容")                
    return
def initSSS(variables, pjName):       ##== 初始化 state_session 的各變量 ==##
    sss = st.session_state
    if "LOG" not in sss:        sss.LOG = [pjName]   #-- 初始化 LOG 列表
    if "prs" not in sss:        sss.prs = Presentation();     sss.prs = genPPTX(pjName,"date: "+str(date.today()))
    if "username" not in sss:   sss.username = ""    #-- 初始化 username 列表
    if "suggestion" not in sss: sss.suggestion = ""  #-- 初始化 suggestion 列表
    for var in variables:       #-- 初始化传入的变量名为 None
        if var not in sss: sss[var] = None
    return sss

#%%##===== (W4).網站架構: (1)設導航列,(2)設sss,(3)設主標題,(4)導航切換,(5)操作日誌 =====#####
if __name__ == "__main__":

    ##== (1).設定頁面組態 與 導航列 (前台(a)navbar) ==##
    st.set_page_config(page_title="AIp08數據關連分析(C01B)", page_icon="✅", layout="wide",)  #==> [[AIp04/C4)(5)加上頁註,頁標題等]]
    page = st_navbar(["[1擷取交易]", "[2轉換客戶圖像]", "[3TFM_客群選取]", "[4目標客群交易]", "[5客群聚類分析]", "[6客群關連規則]", "[7匯出PPT檔]", "(C)"])

    ##== (2).設定session初始值(sss),專案名稱等 ==##
    sss = initSSS(["wkDir","Xname","X","Cv","TFMs","XX","FF0A","MM0A","CvTAtfm"], "聚類數據規C07A"+"--"+wkDir)
    sss.wkDir = wkDir;    sss.Xname = "XXX.csv"

    ##== (3).設定 前台((b)sidebar + (c)canvas)主標題 ==##
    st.title("AIp08數據關連分析(C01B)儀表板")
    st.sidebar.title("控制盤(C01B)--")

    ##== (4).導航切換: 前台(a)navbar-->儀表板函式(b,c,d) ==##
    match page:
        case "[1擷取交易]": sss.X, rv_list, log_list = 擷取交易(sss.wkDir+sss.Xname);    check2log(log_list, sss.LOG, rv_list)
        case "[2轉換客戶圖像]":
            if sss.X is None:    st.write("尚未擷取交易數據，請先擷取交易數據！");    sss.LOG.append("尚未擷取交易數據，請先擷取交易數據！")
            else: sss.Cv, rv_list, log_list = 轉換客戶圖像(sss.X,FFbreaks,MMbreaks,BBbreaks,RRbreaks,Tnow);   check2log(log_list, sss.LOG, rv_list)
        case "[3TFM_客群選取]":
            if sss.Cv is None: st.write("尚未轉換客戶圖像，請先計算客戶圖像！");   sss.LOG.append("尚未轉換客戶圖像，請先轉換客戶圖像！")
            else: sss.TFMs, sss.FF0A, sss.MM0A, sss.CvTAtfm, rv_list, log_list = TFM_客群選取(sss.Cv);      check2log(log_list, sss.LOG, rv_list)            
        case "[4目標客群交易]":
            if sss.CvTAtfm is None: st.write("尚未選取目標客群，請先選取目標客群！");   sss.LOG.append("尚未選取目標客群，請先選取目標客群")
            else: sss.XX, rv_list, log_list = 目標客群交易(sss.CvTAtfm, sss.X);     check2log(log_list, sss.LOG, rv_list)
        case "[5客群聚類分析]":
            if sss.XX is None: st.write("尚未選取目標客群交易，請先選取目標客群！");   sss.LOG.append("尚未選取目標客群交易，請先選取目標客群")
            else: sss.CP_matrix, rv_list, log_list = 客群聚類分析(sss.XX);     check2log(log_list, sss.LOG, rv_list)
        case "[6客群關連規則]":
            if sss.XX is None: st.write("尚未選取目標客群交易，請先選取目標客群！");   sss.LOG.append("尚未選取目標客群交易，請先選取目標客群")
            else: sss.CP_matrix, rv_list, log_list = 客群關連規則(sss.XX);     check2log(log_list, sss.LOG, rv_list)
        case "[7匯出PPT檔]":
            if sss.Cv is None: st.write("尚未轉換客戶圖像，請先計算客戶圖像！");   sss.LOG.append("尚未轉換客戶圖像，請先轉換客戶圖像！")
            else: 匯出PPT檔("聚類數據規畫C01A.PPTX")
        case "(C)":
            st.header("Welcome to JJS AIDA framework (v0.7)")
            st.subheader("Copyright 2024 Joyce Lin and Jia-Sheng Heh. All rights reserved.")
            st.subheader("This software is licensed under the MIT License.")
            
    ##== (5).操作日誌(LOG) ==##
    with st.sidebar.expander("操作日誌(jjs LOG)", expanded=True):    
        st.markdown('<h2 style="color: blue;">操作jjsLOG日誌</h2>', unsafe_allow_html=True)
        for i, log in enumerate(sss.LOG, 1): st.write(f"({i}). {log}")
        if st.sidebar.checkbox("匯出jjsLOG操作日誌(.jjs)"):
            current_time = datetime.now().strftime("%Y%m%d_%H%M")
            file_name = f"流失客分析(C01)_{current_time}.jjs"
            with open(file_name, 'w', encoding='utf-8') as file:
               for entry in sss.LOG: file.write(f"{entry}\n") 
            st.sidebar.write("log已匯出為"+file_name)
